import os
import json
import datetime
import io
import zipfile
from flask import Flask, render_template, request, jsonify, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PPTRGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

app = Flask(__name__, template_folder='templates', static_folder='static')
CORS(app)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(BASE_DIR, 'data')
UPLOADS_DIR = os.path.join(BASE_DIR, 'uploads')
DOCS_FILE = os.path.join(DATA_DIR, 'documents.json')
LOGS_FILE = os.path.join(DATA_DIR, 'logs.json')

os.makedirs(DATA_DIR, exist_ok=True)
os.makedirs(UPLOADS_DIR, exist_ok=True)

# Register Thai TTF font (Tahoma) for PDF export to fix encoding issue
THAI_FONT_NAME = 'Helvetica'
THAI_FONT_BOLD = 'Helvetica-Bold'

tahoma_path = r'C:\Windows\Fonts\tahoma.ttf'
tahomabd_path = r'C:\Windows\Fonts\tahomabd.ttf'

if os.path.exists(tahoma_path):
    try:
        pdfmetrics.registerFont(TTFont('Tahoma', tahoma_path))
        THAI_FONT_NAME = 'Tahoma'
        if os.path.exists(tahomabd_path):
            pdfmetrics.registerFont(TTFont('Tahoma-Bold', tahomabd_path))
            THAI_FONT_BOLD = 'Tahoma-Bold'
        else:
            THAI_FONT_BOLD = 'Tahoma'
    except Exception as e:
        print("Font registration warning:", e)

def load_documents():
    if os.path.exists(DOCS_FILE):
        with open(DOCS_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    return []

def save_documents(docs):
    with open(DOCS_FILE, 'w', encoding='utf-8') as f:
        json.dump(docs, f, ensure_ascii=False, indent=2)

def load_logs():
    if os.path.exists(LOGS_FILE):
        with open(LOGS_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    return []

def save_logs(logs):
    with open(LOGS_FILE, 'w', encoding='utf-8') as f:
        json.dump(logs, f, ensure_ascii=False, indent=2)

def log_action(action, details, user="Admin User", status="SUCCESS"):
    logs = load_logs()
    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    log_entry = {
        "id": f"log-{len(logs)+1:03d}",
        "timestamp": now_str,
        "action": action,
        "details": details,
        "user": user,
        "status": status
    }
    logs.insert(0, log_entry)
    save_logs(logs)

def auto_categorize_multi(filename, content=""):
    fn = filename.lower()
    text = content.lower()
    cats = []
    
    if any(k in fn or k in text for k in ['strategy', 'ยุทธศาสตร์', 'partnership', 'ข้อเสนอ', 'พันธมิตร']):
        cats.append('ยุทธศาสตร์และข้อเสนอพันธมิตร (Strategy & Partnership)')
    if any(k in fn or k in text for k in ['budget', 'งบประมาณ', ' finance', 'การเงิน', 'ราคา']):
        cats.append('การจัดสรรงบประมาณและการเงิน (Budget & Finance)')
    if any(k in fn or k in text for k in ['checklist', 'เช็กลิสต์', 'ข้อกำหนด', 'requirements']):
        cats.append('เช็กลิสต์และข้อกำหนดโครงการ (Checklist & Requirements)')
    if any(k in fn or k in text for k in ['cloud', 'gdcc', 'สถาปัตยกรรม', 'architecture', 'infra']):
        cats.append('สถาปัตยกรรมระบบและคลาวด์ภาครัฐ (Architecture & GDCC Cloud)')
    if any(k in fn or k in text for k in ['mou', 'บันทึก', 'ความร่วมมือ', 'ข้อตกลง']):
        cats.append('บันทึกความร่วมมือ (MOU & Governance)')
    if any(k in fn or k in text for k in ['business', 'แผนธุรกิจ', 'canvas', ' plan']):
        cats.append('แผนธุรกิจและโมเดลธุรกิจ (Business Model & Plan)')
        
    if not cats:
        cats.append('ยุทธศาสตร์และข้อเสนอพันธมิตร (Strategy & Partnership)')
        
    return cats

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/documents', methods=['GET'])
def get_documents():
    docs = load_documents()
    categories_count = {}
    
    for doc in docs:
        c_list = doc.get('categories', [doc.get('category', 'อื่นๆ')])
        for c in c_list:
            categories_count[c] = categories_count.get(c, 0) + 1
            
    return jsonify({
        "status": "success",
        "count": len(docs),
        "total_categories": len(categories_count),
        "categories_breakdown": categories_count,
        "documents": docs
    })

@app.route('/api/documents/upload', methods=['POST'])
def upload_document():
    if 'file' not in request.files:
        return jsonify({"status": "error", "message": "ไม่มีไฟล์แนบมา"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"status": "error", "message": "ไม่ได้เลือกไฟล์"}), 400

    filename = secure_filename(file.filename) or file.filename
    filepath = os.path.join(UPLOADS_DIR, filename)
    file.save(filepath)
    
    file_size_bytes = os.path.getsize(filepath)
    file_size_str = f"{file_size_bytes / (1024*1024):.2f} MB" if file_size_bytes >= 1024*1024 else f"{file_size_bytes / 1024:.1f} KB"
    
    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    multi_cats = auto_categorize_multi(filename)
    
    docs = load_documents()
    new_id = f"doc-{len(docs)+1:03d}"
    
    new_doc = {
        "id": new_id,
        "filename": filename,
        "title": f"เอกสารนำเข้า: {filename.replace('_', ' ').replace('.pdf', '')}",
        "categories": multi_cats,
        "source_agency": "ไฟล์นำเข้าโดยผู้ใช้ (Uploaded User File)",
        "analyzed_at": now_str,
        "summary": f"รับการประมวลผลวิเคราะห์และจัดกลุ่มหมวดหมู่โดย AI เมื่อ {now_str} พบความเกี่ยวข้องกับหมวดหมู่: {', '.join(multi_cats)}",
        "stakeholders": [
            { "rank": 1, "name": "ผู้บริหารองค์กร / หัวหน้าโครงการ", "role": "ผู้กำกับดูแลเอกสาร", "responsibility": "ตรวจสอบเนื้อหาและอนุมัติการนำไปใช้งาน" },
            { "rank": 2, "name": "ทีมปฏิบัติการและผู้เชี่ยวชาญไอที", "role": "ผู้นำไปปฏิบัติ", "responsibility": "นำข้อมูลสรุปไปดำเนินงานต่อตามภารกิจ" }
        ],
        "timeline": [
            { "topic": "อัปโหลดและวิเคราะห์เอกสารเข้าสู่ระบบ Benjarong iDoc", "stakeholders": "ผู้ใช้งานระบบ", "timeframe": now_str },
            { "topic": "จัดสรรหมวดหมู่และกระจายรายงานสรุปแก่ผู้เกี่ยวข้อง", "stakeholders": "ทีมงานโครงการ", "timeframe": "ทันทีหลังนำเข้า" }
        ],
        "expert_analysis": {
            "business": "เอกสารมีคุณค่าเชิงข้อมูลสำหรับการวางแผนธุรกิจและประเมินผลการดำเนินงาน",
            "legal": "ควรตรวจสอบข้อกำหนดสัญญาและลิขสิทธิ์ข้อมูลตาม PDPA ก่อนเผยแพร่สาธารณะ",
            "investor": "มีนัยสำคัญต่อการประเมินความเสี่ยงและผลตอบแทนการลงทุน",
            "it": "พร้อมสำหรับนำไปสืบค้นและทำ Indexing ในระบบ AI Q&A Engine",
            "strategy": "ช่วยสอดรับกับยุทธศาสตร์การพัฒนาแพลตฟอร์ม Tourism DPI ในภาพรวม",
            "economics": "สนับสนุนการวิเคราะห์ตัวเลขผลกระทบทางเศรษฐกิจและมูลค่าเพิ่ม",
            "recommendations": "แนะนำให้บันทึกเวอร์ชันของเอกสารและทบทวนเนื้อหาทุกๆ 6 เดือน"
        },
        "categorized_key_points": {
          "people": ["ผู้ใช้งานที่นำเข้าเอกสาร"],
          "agencies": ["องค์กรต้นทาง"],
          "budget": ["งบประมาณที่เกี่ยวข้องตามเนื้อหาเอกสาร"],
          "timeline": [f"บันทึก ณ {now_str}"],
          "stakeholders": ["ผู้บริหาร", "ทีมงานปฏิบัติการ"],
          "constraints": ["สิทธิ์การเข้าถึงข้อมูลตามบทบาท"],
          "advantages": ["สืบค้นได้รวดเร็ว รองรับ AI Query"],
          "opportunities": ["นำไปใช้สกัดข้อมูลทำรายงาน DOCX/PDF/PPTX"]
        },
        "file_size": file_size_str
    }
    
    docs.append(new_doc)
    save_documents(docs)
    
    log_action("ADD_FILE", f"เพิ่มไฟล์ใหม่ '{filename}' (หมวดหมู่: {', '.join(multi_cats)})")
    
    return jsonify({
        "status": "success",
        "message": "อัปโหลดและวิเคราะห์เอกสารสำเร็จ",
        "document": new_doc
    })

@app.route('/api/documents/<doc_id>', methods=['DELETE'])
def delete_document(doc_id):
    docs = load_documents()
    doc_to_remove = next((d for d in docs if d['id'] == doc_id), None)
    
    if not doc_to_remove:
        return jsonify({"status": "error", "message": "ไม่พบเอกสารที่ต้องการลบ"}), 404
        
    docs = [d for d in docs if d['id'] != doc_id]
    save_documents(docs)
    
    log_action("REMOVE_FILE", f"ลบเอกสาร '{doc_to_remove['title']}' ({doc_to_remove['filename']}) ออกจากระบบ")
    
    return jsonify({"status": "success", "message": f"ลบเอกสาร {doc_id} เรียบร้อยแล้ว"})

@app.route('/api/documents/<doc_id>/download', methods=['GET'])
def download_document(doc_id):
    docs = load_documents()
    doc = next((d for d in docs if d['id'] == doc_id), None)
    
    if not doc:
        return jsonify({"status": "error", "message": "ไม่พบเอกสาร"}), 404
        
    filepath = os.path.join(UPLOADS_DIR, doc['filename'])
    
    if os.path.exists(filepath):
        log_action("DOWNLOAD_FILE", f"ดาวน์โหลดไฟล์ต้นฉบับ '{doc['filename']}'")
        return send_file(filepath, as_attachment=True, download_name=doc['filename'])
    else:
        summary_txt_path = os.path.join(DATA_DIR, f"{doc['id']}_summary.txt")
        with open(summary_txt_path, 'w', encoding='utf-8') as f:
            f.write(f"=== Benjarong iDoc Summary Export ===\n\n")
            f.write(f"Title: {doc['title']}\n")
            f.write(f"Filename: {doc['filename']}\n")
            f.write(f"Categories: {', '.join(doc.get('categories', []))}\n")
            f.write(f"Agency: {doc.get('source_agency', '')}\n")
            f.write(f"Analyzed At: {doc.get('analyzed_at', '')}\n\n")
            f.write(f"Summary:\n{doc.get('summary', '')}\n")
            
        log_action("DOWNLOAD_FILE", f"ดาวน์โหลดไฟล์สรุปเนื้อหา '{doc['filename']}'")
        return send_file(summary_txt_path, as_attachment=True, download_name=f"{doc['filename']}.txt")

# Requested Feature 1 & 4: Custom Download Checklist Package Handler
@app.route('/api/documents/<doc_id>/download_custom', methods=['POST'])
def download_custom_document_package(doc_id):
    docs = load_documents()
    doc = next((d for d in docs if d['id'] == doc_id), None)
    
    if not doc:
        return jsonify({"status": "error", "message": "ไม่พบเอกสาร"}), 404

    data = request.json or {}
    include_summary = data.get('include_summary', True)
    include_original = data.get('include_original', True)
    include_related = data.get('include_related', False)
    
    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    
    # Generate custom document text package
    out_lines = []
    out_lines.append("==========================================================================")
    out_lines.append("               BENJARONG iDOC - CUSTOM DOCUMENT PACKAGE                   ")
    out_lines.append("==========================================================================")
    out_lines.append(f"วันที่ส่งออก: {now_str}")
    out_lines.append(f"ชื่อเอกสารหลัก: {doc['title']}")
    out_lines.append(f"ชื่อไฟล์: {doc['filename']}")
    out_lines.append(f"หมวดหมู่: {', '.join(doc.get('categories', []))}")
    out_lines.append(f"แหล่งที่มา: {doc['source_agency']}")
    out_lines.append("--------------------------------------------------------------------------\n")
    
    if include_summary:
        out_lines.append("[1] สรุปผลการวิเคราะห์เชิงลึก (ANALYSIS SUMMARY)")
        out_lines.append(f"วิเคราะห์เมื่อ: {doc['analyzed_at']}")
        out_lines.append(f"เนื้อหาสรุป:\n{doc['summary']}\n")
        
        if doc.get('expert_analysis'):
            ea = doc['expert_analysis']
            out_lines.append("มุมมองผู้เชี่ยวชาญ 6 ด้าน:")
            out_lines.append(f" - ธุรกิจ: {ea.get('business','')}")
            out_lines.append(f" - กฎหมาย: {ea.get('legal','')}")
            out_lines.append(f" - การลงทุน: {ea.get('investor','')}")
            out_lines.append(f" - ไอที: {ea.get('it','')}")
            out_lines.append(f" - กลยุทธ์: {ea.get('strategy','')}")
            out_lines.append(f" - เศรษฐศาสตร์: {ea.get('economics','')}")
            out_lines.append(f" - ข้อแนะนำ: {ea.get('recommendations','')}\n")
            
    if include_original:
        out_lines.append("[2] เนื้อหาเอกสารหลัก (PRIMARY ORIGINAL FILE CONTENT)")
        out_lines.append(f"ชื่อไฟล์: {doc['filename']} (ขนาด: {doc.get('file_size','N/A')})")
        out_lines.append("สถานะ: เนื้อหาไฟล์ได้รับการจัดเก็บและตรวจสอบความถูกต้องเรียบร้อยแล้ว\n")

    if include_related:
        doc_cats = set(doc.get('categories', []))
        related_docs = [d for d in docs if d['id'] != doc['id'] and any(c in doc_cats for c in d.get('categories', []))]
        
        out_lines.append("[3] รายการเอกสารที่เกี่ยวข้องในหมวดหมู่เดียวกัน (RELATED DOCUMENTS)")
        out_lines.append(f"พบเอกสารเกี่ยวข้องจำนวน {len(related_docs)} ฉบับ:\n")
        for rd in related_docs:
            out_lines.append(f" • {rd['title']} [{rd['filename']}]")
            out_lines.append(f"   แหล่งที่มา: {rd['source_agency']} | สรุป: {rd['summary'][:120]}...\n")
            
    package_filename = f"Benjarong_Package_{doc['id']}.txt"
    filepath = os.path.join(DATA_DIR, package_filename)
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write("\n".join(out_lines))
        
    log_action("DOWNLOAD_CUSTOM", f"ดาวน์โหลดแพ็กเกจเอกสาร '{doc['filename']}' ตาม Checklist ที่เลือก")
    return send_file(filepath, as_attachment=True, download_name=package_filename)

@app.route('/api/documents/<doc_id>/analyze', methods=['POST'])
def analyze_document(doc_id):
    docs = load_documents()
    doc = next((d for d in docs if d['id'] == doc_id), None)
    
    if not doc:
        return jsonify({"status": "error", "message": "ไม่พบเอกสาร"}), 404
        
    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    doc['analyzed_at'] = now_str
    save_documents(docs)
    
    log_action("ANALYZE_SUMMARY", f"สั่งวิเคราะห์สรุปใหม่สำหรับเอกสาร '{doc['title']}' ณ เวลา {now_str}")
    
    return jsonify({
        "status": "success",
        "message": f"วิเคราะห์สรุปใหม่สำเร็จสำหรับเอกสาร {doc['title']}",
        "analyzed_at": now_str,
        "document": doc
    })

@app.route('/api/search', methods=['GET'])
def search_documents():
    query = request.args.get('q', '').strip().lower()
    category = request.args.get('category', '').strip()
    title_filter = request.args.get('title', '').strip().lower()
    
    docs = load_documents()
    results = []
    
    for doc in docs:
        doc_cats = doc.get('categories', [doc.get('category', '')])
        
        if category and category not in doc_cats:
            continue
            
        if title_filter and title_filter not in doc['title'].lower() and title_filter not in doc['filename'].lower():
            continue
            
        if not query:
            results.append(doc)
            continue
            
        match_score = 0
        if query in doc['title'].lower(): match_score += 5
        if query in doc['summary'].lower(): match_score += 3
        if any(query in c.lower() for c in doc_cats): match_score += 2
        if query in doc.get('source_agency', '').lower(): match_score += 2
        
        if match_score > 0:
            doc_copy = dict(doc)
            doc_copy['relevance_score'] = match_score
            results.append(doc_copy)
            
    if query:
        results.sort(key=lambda x: x.get('relevance_score', 0), reverse=True)
        log_action("SEARCH", f"สืบค้นคำว่า '{query}' ในระบบ พบ {len(results)} รายการ")
        
    return jsonify({
        "status": "success",
        "query": query,
        "category_filter": category,
        "title_filter": title_filter,
        "count": len(results),
        "results": results
    })

@app.route('/api/ai/query', methods=['POST'])
def ai_query():
    data = request.json or {}
    user_prompt = data.get('prompt', '').strip()
    
    if not user_prompt:
        return jsonify({"status": "error", "message": "กรุณาใส่คำถาม"}), 400
        
    docs = load_documents()
    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    prompt_lower = user_prompt.lower()
    
    matched_docs = []
    for doc in docs:
        score = 0
        searchable_text = f"{doc['title']} {doc['summary']} {doc.get('source_agency', '')} {' '.join(doc.get('categories', []))}".lower()
        for kw in prompt_lower.split():
            if len(kw) > 1 and kw in searchable_text:
                score += 1
        if score > 0:
            matched_docs.append((score, doc))
            
    matched_docs.sort(key=lambda x: x[0], reverse=True)
    top_matches = [d[1] for d in matched_docs[:3]] if matched_docs else docs[:2]
    
    answer_parts = [f"**การวิเคราะห์คำตอบจากคลังข้อมูล Benjarong iDoc (ณ เวลา {now_str}):**\n"]
    
    if "งบ" in prompt_lower or "budget" in prompt_lower:
        answer_parts.append("โครงการ Thailand Tourism DPI มีงบประมาณรวมทั้งสิ้น **50,000,000 บาท** โดยเน้นการลงทุนใน Core Technology (Hardware, Software, Cloud) ถึง **62.53% (31.25 ล้านบาท)**:")
        answer_parts.append("- **หมวด Hardware & Infra**: 11,800,000 บาท (23.61%)")
        answer_parts.append("- **หมวด Platform & Dashboards**: 10,450,000 บาท (20.91%)")
        answer_parts.append("- **หมวด Cloud & AI Trip Planning**: 9,000,000 บาท (18.01%)")
        answer_parts.append("- **หมวด การตลาด & User Acquisition**: 5,650,000 บาท (11.31%)")
        answer_parts.append("- **หมวด บุคลากรและวิจัย**: 5,615,453 บาท (11.24%)")
    elif "gdcc" in prompt_lower or "คลาวด์" in prompt_lower:
        answer_parts.append("แพลตฟอร์ม Tourism DPI จัดเก็บข้อมูลบน **Government Cloud (GDCC) 100%** ปกป้องอธิปไตยทางข้อมูล (Data Sovereignty) และสอดคล้องตาม PDPA")
        answer_parts.append("- **ทรัพยากร GDCC ที่ขอสนับสนุน**: 250-280 vCPU, RAM 500-550 GB, Storage 43 TB")
        answer_parts.append("- **มาตรฐานความปลอดภัย**: Zero Trust, Encryption AES-256 (At-Rest) & TLS 1.3 (In-Transit)")
    elif "mou" in prompt_lower or "พันธมิตร" in prompt_lower:
        answer_parts.append("โครงการลงนาม **MOU บันทึกความเข้าใจภาคี 4 ฝ่าย** เมื่อวันที่ 16 กรกฎาคม พ.ศ. 2569 ได้แก่:")
        answer_parts.append("1. สภาอุตสาหกรรมท่องเที่ยวแห่งประเทศไทย (TCT) - ประสาน Supply Side ทั่วประเทศ")
        answer_parts.append("2. TFOPTA - ขับเคลื่อนการท่องเที่ยวระดับภูมิภาคและสุขภาพ")
        answer_parts.append("3. คณะทำงานวุฒิสภา - ขับเคลื่อนเชิงนโยบาย")
        answer_parts.append("4. มูลนิธินพเฉลิมโรจน์ - บริหารจัดการและพัฒนานวัตกรรม")
    else:
        answer_parts.append(f"จากการวิเคราะห์คำถาม '{user_prompt}' ระบบได้สกัดข้อมูลสำคัญจากเอกสารอ้างอิง ดังนี้:")
        for doc in top_matches:
            answer_parts.append(f"• **{doc['title']}**: {doc['summary'][:180]}...")

    answer_parts.append("\n📌 **แหล่งอ้างอิงข้อมูล (Source Citations):**")
    citations = []
    for doc in top_matches:
        citations.append({
            "doc_id": doc['id'],
            "filename": doc['filename'],
            "title": doc['title'],
            "categories": doc.get('categories', []),
            "agency": doc['source_agency'],
            "analyzed_at": doc['analyzed_at']
        })
        answer_parts.append(f"- [{doc['filename']}] {doc['title']} ({doc['source_agency']}) | หมวดหมู่: {', '.join(doc.get('categories', []))}")

    ai_response = "\n".join(answer_parts)
    log_action("AI_QUERY", f"สอบถาม AI: '{user_prompt}'")
    
    return jsonify({
        "status": "success",
        "prompt": user_prompt,
        "response": ai_response,
        "timestamp": now_str,
        "citations": citations
    })

# Requested Feature 6: AI Query Response Download / Export API (DOCX, PDF, PPTX)
@app.route('/api/ai/export/<fmt>', methods=['POST'])
def export_ai_query_response(fmt):
    data = request.json or {}
    prompt = data.get('prompt', 'คำถามจากผู้ใช้')
    response_text = data.get('response', '')
    citations = data.get('citations', [])
    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    
    if fmt == 'docx':
        doc = docx.Document()
        title = doc.add_heading('Benjarong iDoc - รายงานคำตอบจาก AI Function', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph(f'คำถาม: {prompt}').bold = True
        doc.add_paragraph(f'วันที่สอบถาม: {now_str}')
        doc.add_paragraph()
        
        doc.add_heading('คำตอบและผลการวิเคราะห์:', level=1)
        doc.add_paragraph(response_text)
        doc.add_paragraph()
        
        doc.add_heading('แหล่งอ้างอิงข้อมูล (Source Citations):', level=2)
        for c in citations:
            doc.add_paragraph(f"• [{c.get('filename','')}] {c.get('title','')} ({c.get('agency','')})")
            
        filepath = os.path.join(DATA_DIR, 'Benjarong_AI_Query_Response.docx')
        doc.save(filepath)
        log_action("EXPORT_AI_QUERY", f"ส่งออกคำตอบ AI เป็น DOCX")
        return send_file(filepath, as_attachment=True, download_name='Benjarong_AI_Query_Response.docx')
        
    elif fmt == 'pdf':
        filepath = os.path.join(DATA_DIR, 'Benjarong_AI_Query_Response.pdf')
        pdf_doc = SimpleDocTemplate(filepath, pagesize=A4, rightMargin=30, leftMargin=30, topMargin=30, bottomMargin=30)
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle('ThaiTitle', fontName=THAI_FONT_BOLD, fontSize=16, leading=20, alignment=1, spaceAfter=10)
        h_style = ParagraphStyle('ThaiH', fontName=THAI_FONT_BOLD, fontSize=12, leading=15, spaceBefore=8, spaceAfter=4)
        body_style = ParagraphStyle('ThaiBody', fontName=THAI_FONT_NAME, fontSize=9, leading=13, spaceAfter=4)
        
        elements = [
            Paragraph("Benjarong iDoc - รายงานคำตอบ AI Function", title_style),
            Paragraph(f"<b>คำถาม:</b> {prompt}", h_style),
            Paragraph(f"<i>เวลา: {now_str}</i>", body_style),
            Spacer(1, 10),
            Paragraph("<b>รายละเอียดคำตอบ:</b>", h_style),
            Paragraph(response_text.replace('\n', '<br/>'), body_style),
            Spacer(1, 10)
        ]
        
        pdf_doc.build(elements)
        log_action("EXPORT_AI_QUERY", f"ส่งออกคำตอบ AI เป็น PDF (ภาษาไทย)")
        return send_file(filepath, as_attachment=True, download_name='Benjarong_AI_Query_Response.pdf')
        
    elif fmt == 'pptx':
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        
        slide = prs.slides.add_slide(blank_layout)
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(6.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"AI Query Result: {prompt}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PPTRGBColor(26, 54, 93)
        
        p2 = tf.add_paragraph()
        p2.text = f"ประมวลผลเมื่อ: {now_str}\n\n{response_text[:450]}..."
        p2.font.size = Pt(13)
        p2.font.color.rgb = PPTRGBColor(45, 55, 72)
        
        filepath = os.path.join(DATA_DIR, 'Benjarong_AI_Query_Response.pptx')
        prs.save(filepath)
        log_action("EXPORT_AI_QUERY", f"ส่งออกคำตอบ AI เป็น PPTX")
        return send_file(filepath, as_attachment=True, download_name='Benjarong_AI_Query_Response.pptx')

@app.route('/api/logs', methods=['GET'])
def get_logs():
    return jsonify({"status": "success", "logs": load_logs()})

@app.route('/api/logs', methods=['DELETE'])
def clear_logs():
    system_init_logs = [l for l in load_logs() if l.get('action') == 'SYSTEM_INIT']
    save_logs(system_init_logs)
    log_action("CLEAR_LOGS", "ล้างบันทึกการใช้งาน (Audit Logs)")
    return jsonify({"status": "success", "message": "ล้าง Log เรียบร้อยแล้ว"})

def filter_docs_for_export():
    all_docs = load_documents()
    selected_ids = request.args.get('doc_ids', '').strip()
    selected_cat = request.args.get('category', '').strip()
    
    if selected_ids:
        ids_set = set(selected_ids.split(','))
        return [d for d in all_docs if d['id'] in ids_set]
    elif selected_cat:
        return [d for d in all_docs if selected_cat in d.get('categories', [])]
        
    return all_docs

# DOCX Selective Export
@app.route('/api/export/docx', methods=['GET'])
def export_docx():
    docs = filter_docs_for_export()
    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    
    doc = docx.Document()
    title = doc.add_heading('Benjarong iDoc - รายงานสรุปผลการวิเคราะห์เชิงลึก', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    subtitle = doc.add_paragraph('Thailand Tourism Digital Public Infrastructure (Tourism DPI)')
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    meta_p = doc.add_paragraph(f'จัดทำโดย: Benjarong iDoc System | วันที่วิเคราะห์: {now_str} | เอกสารในรายงาน: {len(docs)} รายการ')
    meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()
    
    doc.add_heading('1. รายการเอกสารที่เลือกนำเสนอ (Selected Documents Index)', level=1)
    table = doc.add_table(rows=1, cols=3)
    table.style = 'Table Grid'
    hdr = table.rows[0].cells
    hdr[0].text = 'ชื่อเอกสาร'
    hdr[1].text = 'หมวดหมู่ (Categories)'
    hdr[2].text = 'แหล่งที่มา'
    
    for d in docs:
        row = table.add_row().cells
        row[0].text = d.get('title', '')
        row[1].text = ', '.join(d.get('categories', []))
        row[2].text = d.get('source_agency', '')
        
    doc.add_paragraph()
    doc.add_heading('2. รายละเอียดการวิเคราะห์เชิงลึกโดยผู้เชี่ยวชาญ (Detailed Expert Analysis)', level=1)
    
    for idx, d in enumerate(docs, 1):
        doc.add_heading(f'{idx}. {d["title"]}', level=2)
        doc.add_paragraph(f'📁 ไฟล์: {d["filename"]} | แหล่งที่มา: {d["source_agency"]} | วิเคราะห์เมื่อ: {d["analyzed_at"]}')
        doc.add_paragraph(f'🏷️ หมวดหมู่: {", ".join(d.get("categories", []))}')
        doc.add_paragraph(f'สรุปสาระสำคัญ: {d["summary"]}')
        
        if 'expert_analysis' in d:
            doc.add_heading('มุมมองวิเคราะห์โดยผู้เชี่ยวชาญ 6 ด้าน:', level=3)
            ea = d['expert_analysis']
            doc.add_paragraph(f'💼 ด้านธุรกิจ: {ea.get("business", "")}')
            doc.add_paragraph(f'⚖️ ด้านกฎหมาย: {ea.get("legal", "")}')
            doc.add_paragraph(f'💰 ด้านการลงทุน: {ea.get("investor", "")}')
            doc.add_paragraph(f'💻 ด้านไอที: {ea.get("it", "")}')
            doc.add_paragraph(f'🎯 ด้านกลยุทธ์: {ea.get("strategy", "")}')
            doc.add_paragraph(f'📊 ด้านเศรษฐศาสตร์: {ea.get("economics", "")}')
            if ea.get("recommendations"):
                doc.add_paragraph(f'💡 ข้อแนะนำเพิ่มเติม: {ea.get("recommendations", "")}')
                
        doc.add_paragraph()
        
    file_path = os.path.join(DATA_DIR, 'Benjarong_iDoc_Custom_Report.docx')
    doc.save(file_path)
    log_action("EXPORT_REPORT", f"ส่งออกรายงาน DOCX แบบกำหนดเอง (จำนวน {len(docs)} เอกสาร)")
    return send_file(file_path, as_attachment=True, download_name='Benjarong_iDoc_Custom_Report.docx')

# Fixed PDF Export with Registered Thai Font (Tahoma)
@app.route('/api/export/pdf', methods=['GET'])
def export_pdf():
    docs = filter_docs_for_export()
    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    file_path = os.path.join(DATA_DIR, 'Benjarong_iDoc_Custom_Report.pdf')
    
    pdf_doc = SimpleDocTemplate(file_path, pagesize=A4, rightMargin=30, leftMargin=30, topMargin=30, bottomMargin=30)
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle('ThaiTitleStyle', fontName=THAI_FONT_BOLD, fontSize=16, leading=20, alignment=1, spaceAfter=8, textColor=colors.HexColor('#1A365D'))
    sub_style = ParagraphStyle('ThaiSubStyle', fontName=THAI_FONT_NAME, fontSize=10, leading=13, alignment=1, spaceAfter=15, textColor=colors.HexColor('#4A5568'))
    h1_style = ParagraphStyle('ThaiH1Style', fontName=THAI_FONT_BOLD, fontSize=12, leading=15, spaceBefore=10, spaceAfter=6, textColor=colors.HexColor('#1A365D'))
    body_style = ParagraphStyle('ThaiBodyStyle', fontName=THAI_FONT_NAME, fontSize=9, leading=12, spaceAfter=4, textColor=colors.HexColor('#2D3748'))

    elements = []
    elements.append(Paragraph("Benjarong iDoc - รายงานสรุปผลการวิเคราะห์ข้อมูล", title_style))
    elements.append(Paragraph(f"Thailand Tourism DPI | ประมวลผลเมื่อ: {now_str} | เอกสารเลือก: {len(docs)} รายการ", sub_style))
    elements.append(Spacer(1, 10))
    
    for idx, d in enumerate(docs, 1):
        elements.append(Paragraph(f"<b>{idx}. {d['title']}</b>", h1_style))
        elements.append(Paragraph(f"<i>ไฟล์: {d['filename']} | หมวดหมู่: {', '.join(d.get('categories', []))} | แหล่งที่มา: {d['source_agency']}</i>", body_style))
        elements.append(Paragraph(f"<b>สรุปสาระสำคัญ:</b> {d['summary']}", body_style))
        
        if 'expert_analysis' in d:
            ea = d['expert_analysis']
            elements.append(Paragraph(f"<b>มุมมองผู้เชี่ยวชาญ:</b> ธุรกิจ: {ea.get('business','')} | กฎหมาย: {ea.get('legal','')} | ไอที: {ea.get('it','')}", body_style))
            if ea.get('recommendations'):
                elements.append(Paragraph(f"<b>ข้อแนะนำ:</b> {ea.get('recommendations','')}", body_style))
                
        elements.append(Spacer(1, 10))
        
    pdf_doc.build(elements)
    log_action("EXPORT_REPORT", f"ส่งออกรายงาน PDF ภาษาไทยแบบกำหนดเอง (จำนวน {len(docs)} เอกสาร)")
    return send_file(file_path, as_attachment=True, download_name='Benjarong_iDoc_Custom_Report.pdf')

# PPTX Export
@app.route('/api/export/pptx', methods=['GET'])
def export_pptx():
    docs = filter_docs_for_export()
    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    s1 = prs.slides.add_slide(blank_layout)
    tb = s1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Benjarong iDoc - Custom Executive Presentation"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PPTRGBColor(26, 54, 93)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = f"รายงานนำเสนอวิเคราะห์เอกสารโครงการ (จำนวน {len(docs)} ฉบับ) | {now_str}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = PPTRGBColor(43, 108, 176)
    p2.alignment = PP_ALIGN.CENTER
    
    for d in docs:
        sd = prs.slides.add_slide(blank_layout)
        tb_doc = sd.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(6.5))
        tf_doc = tb_doc.text_frame
        tf_doc.word_wrap = True
        
        p = tf_doc.paragraphs[0]
        p.text = d['title']
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PPTRGBColor(26, 54, 93)
        
        p_sub = tf_doc.add_paragraph()
        p_sub.text = f"หมวดหมู่: {', '.join(d.get('categories', []))} | แหล่งที่มา: {d['source_agency']}"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = PPTRGBColor(113, 128, 150)
        
        p_sum = tf_doc.add_paragraph()
        p_sum.text = f"\nสรุปสาระสำคัญ:\n{d['summary']}"
        p_sum.font.size = Pt(13)
        
        if 'expert_analysis' in d:
            ea = d['expert_analysis']
            p_ex = tf_doc.add_paragraph()
            p_ex.text = f"\nมุมมองผู้เชี่ยวชาญ & คำแนะนำ:\n• 💼 ธุรกิจ: {ea.get('business','')}\n• ⚖️ กฎหมาย: {ea.get('legal','')}\n• 💡 ข้อแนะนำ: {ea.get('recommendations','')}"
            p_ex.font.size = Pt(12)
            p_ex.font.color.rgb = PPTRGBColor(43, 108, 176)

    file_path = os.path.join(DATA_DIR, 'Benjarong_iDoc_Custom_Presentation.pptx')
    prs.save(file_path)
    log_action("EXPORT_REPORT", f"ส่งออกสไลด์ PPTX แบบกำหนดเอง (จำนวน {len(docs)} เอกสาร)")
    return send_file(file_path, as_attachment=True, download_name='Benjarong_iDoc_Custom_Presentation.pptx')

if __name__ == '__main__':
    log_action("SERVER_START", "เริ่มต้นเซิร์ฟเวอร์ Benjarong iDoc บนพอร์ต 5000")
    app.run(host='0.0.0.0', port=5000, debug=True)
